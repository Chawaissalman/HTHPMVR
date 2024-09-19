# -*- coding: utf-8 -*-
"""
Created on Sun Nov 12 19:04:26 2023

@author: me
"""
import pythoncom  # This is required for COM initialization
from pptx import Presentation

def print_shape_info(ppt_file_path):
    ppt = Presentation(ppt_file_path)
    
    for slide_number, slide in enumerate(ppt.slides):
        print(f"Slide {slide_number + 1}:")
        
        for shape_number, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type
            text = shape.text if shape.has_text_frame else "No text"
            print(f"  Shape {shape_number}: Type={shape_type}, Text='{text}'")

# # Path to the PowerPoint file
# ppt_path = r'C:\Users\me\OneDrive\001 siemens\models\pyhton\HTHP\HTHP_HP+MMVR.pptx'

# # Print shape information for each slide
# print_shape_info(ppt_path)

from pptx import Presentation
from pptx.util import Pt  # Import the Pt class for specifying font size in points

def update_text_while_preserving_formatting(ppt_file_path, updates):
    ppt = Presentation(ppt_file_path)
    
    for update in updates:
        slide = ppt.slides[update['slide_number']]
        shape = slide.shapes[update['shape_number']]
        
        if not shape.has_text_frame:
            print(f"No text frame found in slide {update['slide_number']}, shape {update['shape_number']}")
            continue
        
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        
        # Clear the existing text in the first paragraph
        p.clear()
        # Create a new run with the new text
        run = p.add_run()
        run.text = update['new_text']

        # Set the font properties to match the original formatting
        font = run.font
        font.name = 'Arial'
        font.size = Pt(10)
        
        
    
    # Save the updated presentation
    updated_ppt_path = ppt_file_path.replace('.pptx', '_updated.pptx')
    ppt.save(updated_ppt_path)
    return updated_ppt_path


from pptx import Presentation

def replace_specific_text(ppt_file_path, updates):
    ppt = Presentation(ppt_file_path)

    for update in updates:
        slide = ppt.slides[update['slide_number']]
        shape = slide.shapes[update['shape_number']]

        if not shape.has_text_frame:
            continue

        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if 'xx' in run.text:
                    run.text = run.text.replace('xx', update['new_text'])

    updated_ppt_path = ppt_file_path.replace('.pptx', '_replaced.pptx')
    ppt.save(updated_ppt_path)
    return updated_ppt_path

'-----------------------------------------------------------------------------------------------'

import streamlit as st
from pptx import Presentation
from io import BytesIO
from pdf2image import convert_from_path
import tempfile

import win32com.client

def save_slide_as_image(ppt_path, slide_number, output_path):
                try:
                    pythoncom.CoInitialize()  # Initialize COM library
                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    powerpoint.Visible = 1
                    presentation = powerpoint.Presentations.Open(ppt_path)
                    slide = presentation.Slides[slide_number]
                    slide.Export(output_path, "PNG")
                    presentation.Close()
                    powerpoint.Quit()
                except Exception as e:
                    print(f"An error occurred: {e}")
                    if powerpoint:
                        powerpoint.Quit()
                finally:
                    pythoncom.CoUninitialize()  # Uninitialize COM library


